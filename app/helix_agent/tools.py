import os
import glob
import duckdb
from jinja2 import Environment, FileSystemLoader
from xhtml2pdf import pisa
from pptx import Presentation

# --- DUCKDB TOOL ---
def execute_duckdb_query(query: str) -> str:
    """Executes a SQL query using DuckDB and returns the result as a string."""
    try:
        result = duckdb.query(query).df()
        return result.to_string()
    except Exception as e:
        return f"Query Failed: {str(e)}. Please correct your SQL."

# --- SKILL TOOLS ---
def get_skills_dir() -> str:
    # Look for the 'skills' folder in the same directory as this file
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), "skills")

def list_skills() -> str:
    """Lists all available markdown skills in the skills directory."""
    skills_dir = get_skills_dir()
    if not os.path.exists(skills_dir):
        return "No skills directory found."
    
    skills = []
    for filepath in glob.glob(os.path.join(skills_dir, "*.md")):
        skills.append(os.path.basename(filepath))
    
    if not skills:
        return "No skills found."
        
    return "Available skills: " + ", ".join(skills)

def read_skill(skill_name: str) -> str:
    """Reads and returns the contents of a specific skill markdown file. Ensure you pass the correct filename (e.g. 'regulator_perspective.md')."""
    skills_dir = get_skills_dir()
    safe_name = os.path.basename(skill_name)
    if not safe_name.endswith('.md'):
        safe_name += '.md'
        
    skill_path = os.path.join(skills_dir, safe_name)
    if not os.path.exists(skill_path):
        return f"Error: Skill '{safe_name}' not found. Please use list_skills to see available skills."
        
    try:
        with open(skill_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        return f"Error reading skill: {str(e)}"

def generate_pdf_report(design_name: str, report_data: dict, output_filename: str) -> str:
    """Generates a PDF report using HTML templates."""
    try:
        # Sanitize inputs to prevent path traversal
        safe_design = os.path.basename(design_name)
        safe_filename = os.path.basename(output_filename)
        
        template_dir = os.path.join("data", "designs", safe_design)
        if not os.path.exists(template_dir):
            return f"Error: Design '{safe_design}' not found."
            
        env = Environment(loader=FileSystemLoader(template_dir))
        template_name = os.path.basename(report_data.get("template_name", "template.html"))
        template = env.get_template(template_name)
        html_content = template.render(**report_data)
        
        output_path = os.path.join("files", f"{safe_filename}.pdf")
        os.makedirs("files", exist_ok=True)
        
        with open(output_path, "w+b") as result_file:
            pisa_status = pisa.CreatePDF(html_content, dest=result_file)
            
        if pisa_status.err:
            return "Error: Failed to generate PDF."
        return f"Successfully generated PDF at {output_path}"
    except Exception as e:
        return f"Error: {str(e)}"

def generate_ppt_report(design_name: str, report_data: dict, output_filename: str) -> str:
    """Generates a PPTX report using a template."""
    try:
        # Sanitize inputs to prevent path traversal
        safe_design = os.path.basename(design_name)
        safe_filename = os.path.basename(output_filename)
        template_name = os.path.basename(report_data.get("template_name", "template.pptx"))
        
        template_path = os.path.join("data", "designs", safe_design, template_name)
        if not os.path.exists(template_path):
            return f"Error: Design template '{safe_design}/{template_name}' not found."
            
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Dynamic replacement logic
                for key, value in report_data.items():
                    if not isinstance(value, str):
                        continue
                    if f"{{{key}}}" in shape.text:
                        shape.text = shape.text.replace(f"{{{key}}}", value)
                    elif key == "title" and "Title Placeholder" in shape.text:
                        shape.text = value
                    elif key == "body" and "Subtitle Placeholder" in shape.text:
                        shape.text = value
                    
        output_path = os.path.join("files", f"{safe_filename}.pptx")
        os.makedirs("files", exist_ok=True)
        prs.save(output_path)
        
        return f"Successfully generated PPTX at {output_path}"
    except Exception as e:
        return f"Error: {str(e)}"

# --- REGISTRY ---
REGISTRY = {
    "execute_duckdb_query": execute_duckdb_query,
    "list_skills": list_skills,
    "read_skill": read_skill,
    "generate_pdf_report": generate_pdf_report,
    "generate_ppt_report": generate_ppt_report,
}
