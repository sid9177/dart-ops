from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor

def main():
    BASE_DIR = Path(__file__).resolve().parent.parent
    OUTPUT_DIR = BASE_DIR / "data" / "designs" / "executive_summary"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Citi Title Placeholder"
    if title.text_frame.paragraphs:
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x3B, 0x70)
        
    subtitle.text = "Citi Subtitle Placeholder"
    if subtitle.text_frame.paragraphs:
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x3B, 0x70)
        
    output_path = OUTPUT_DIR / "template.pptx"
    prs.save(str(output_path))

if __name__ == '__main__':
    main()
