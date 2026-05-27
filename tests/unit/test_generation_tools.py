import os
import shutil
import pytest
from app.helix_agent.tools import generate_pdf_report, generate_ppt_report
from pptx import Presentation

# Setup mock data directories for testing
@pytest.fixture
def setup_test_env(tmp_path):
    # Setup data/designs mock
    designs_dir = tmp_path / "data" / "designs" / "mock_design"
    designs_dir.mkdir(parents=True, exist_ok=True)
    
    # Create mock HTML template
    html_template = designs_dir / "template.html"
    html_template.write_text("<html><body><h1>{{ title }}</h1><p>{{ body }}</p></body></html>")
    
    # Create mock PPTX template
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Title Placeholder"
    subtitle = slide.placeholders[1]
    subtitle.text = "{body}"
    prs.save(designs_dir / "template.pptx")
    
    # Change working directory to tmp_path during test
    old_cwd = os.getcwd()
    os.chdir(tmp_path)
    yield tmp_path
    os.chdir(old_cwd)

def test_generate_pdf_report_success(setup_test_env):
    report_data = {"title": "Test PDF", "body": "This is a test PDF body."}
    result = generate_pdf_report("mock_design", report_data, "test_output")
    
    assert "Successfully generated PDF" in result
    assert os.path.exists(os.path.join("files", "test_output.pdf"))

def test_generate_pdf_report_missing_design(setup_test_env):
    result = generate_pdf_report("nonexistent_design", {}, "test_output")
    assert "Error: Design 'nonexistent_design' not found." in result

def test_generate_ppt_report_success(setup_test_env):
    report_data = {"title": "Test PPT", "body": "Replaced body text."}
    result = generate_ppt_report("mock_design", report_data, "test_output")
    
    assert "Successfully generated PPTX" in result
    output_path = os.path.join("files", "test_output.pptx")
    assert os.path.exists(output_path)
    
    # Verify content replacement
    prs = Presentation(output_path)
    slide = prs.slides[0]
    texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
    assert any("Test PPT" in text for text in texts)
    assert any("Replaced body text." in text for text in texts)

def test_generate_ppt_report_missing_design(setup_test_env):
    result = generate_ppt_report("nonexistent_design", {}, "test_output")
    assert "Error: Design template 'nonexistent_design/template.pptx' not found." in result

def test_path_traversal_prevention(setup_test_env):
    report_data = {"title": "Test", "body": "Test"}
    
    # Test path traversal in design name
    result = generate_pdf_report("../../../etc/passwd", report_data, "test_output")
    assert "Error: Design 'passwd' not found." in result
    
    # Test path traversal in output filename
    result = generate_pdf_report("mock_design", report_data, "../../../etc/evil")
    assert "Successfully generated PDF" in result
    assert os.path.exists(os.path.join("files", "evil.pdf"))

def test_unhandled_exception_handling(setup_test_env):
    # Force an exception by providing invalid report data type (e.g., passing a string instead of dict)
    # The render function expects dict unpacking. Wait, the ** operator requires a mapping.
    # Passing a string to ** will raise TypeError before function call, which we can't catch inside.
    # Instead, let's provide report_data that causes jinja to fail or similar, or just test PPTX with broken pptx
    
    # Break the PPTX file to cause an OSError
    with open(os.path.join("data", "designs", "mock_design", "template.pptx"), "w") as f:
        f.write("corrupted data")
        
    result = generate_ppt_report("mock_design", {}, "test_output")
    assert "Error:" in result
    assert "Successfully" not in result
