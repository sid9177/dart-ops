import sys
import io
import os
import traceback
import html
import subprocess

def execute_python_code(code: str) -> str:
    wrapper = (
        "import sys\n"
        "import traceback\n"
        "exec_globals = {}\n"
        "try:\n"
        f"    exec({repr(code)}, exec_globals)\n"
        "except BaseException:\n"
        "    sys.stderr.write('Execution Error:\\n' + traceback.format_exc())\n"
        "    sys.exit(1)\n"
    )
    try:
        res = subprocess.run(
            [sys.executable, "-c", wrapper],
            capture_output=True,
            text=True,
            timeout=10
        )
        return res.stdout + res.stderr
    except subprocess.TimeoutExpired:
        return "Execution Error: Code execution timed out (exceeded 10 seconds)."
    except Exception as e:
        return f"Execution Error: {str(e)}"

def parse_markdown(markdown_content: str):
    elements = []
    current_paragraph = []

    def flush_paragraph():
        if current_paragraph:
            elements.append({
                "type": "paragraph",
                "text": " ".join(current_paragraph)
            })
            current_paragraph.clear()

    lines = markdown_content.split("\n")
    for line in lines:
        stripped_line = line.strip()
        
        # Empty line separates blocks
        if not stripped_line:
            flush_paragraph()
            continue
            
        # Headers
        if stripped_line.startswith("# ") or stripped_line.startswith("## ") or stripped_line.startswith("### "):
            flush_paragraph()
            if stripped_line.startswith("# "):
                level = 1
                text = stripped_line[2:]
            elif stripped_line.startswith("## "):
                level = 2
                text = stripped_line[3:]
            else:
                level = 3
                text = stripped_line[4:]
            elements.append({
                "type": "header",
                "level": level,
                "text": text
            })
            continue

        # Bullet points: e.g. "  * subbullet" or "* bullet"
        leading_spaces = len(line) - len(line.lstrip(' '))
        if stripped_line.startswith("* ") or stripped_line.startswith("- "):
            flush_paragraph()
            level = leading_spaces // 2
            text = stripped_line[2:]
            elements.append({
                "type": "bullet",
                "level": level,
                "text": text
            })
            continue

        # Normal text line - part of paragraph
        current_paragraph.append(stripped_line)

    flush_paragraph()
    return elements

def export_report_to_pdf(markdown_content: str, output_path: str):
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    except ImportError:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"PDF FALLBACK - MOCK GENERATED\n\n{markdown_content}")
        return

    doc = SimpleDocTemplate(output_path, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    
    citi_blue = colors.HexColor("#003B70")
    citi_red = colors.HexColor("#EE3124")
    charcoal = colors.HexColor("#222222")

    title_style = ParagraphStyle(
        'CitiTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        textColor=citi_blue,
        spaceAfter=15
    )
    
    h2_style = ParagraphStyle(
        'CitiH2',
        parent=title_style,
        fontSize=16,
        spaceBefore=10
    )
    
    h3_style = ParagraphStyle(
        'CitiH3',
        parent=title_style,
        fontSize=12,
        spaceBefore=8
    )
    
    body_style = ParagraphStyle(
        'CitiBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10.5,
        textColor=charcoal,
        spaceAfter=8,
        leading=14
    )

    # Dynamic style helper for nested bullets
    def get_bullet_style(level):
        style_name = f'CitiBulletL{level}'
        if style_name in styles:
            return styles[style_name]
        indent = 15 + level * 15
        new_style = ParagraphStyle(
            style_name,
            parent=body_style,
            leftIndent=indent,
            firstLineIndent=-10
        )
        styles.add(new_style)
        return new_style

    story = []
    
    header_bar = Table([[""]], colWidths=[540])
    header_bar.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), citi_red),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('TOPPADDING', (0,0), (-1,-1), 0),
    ]))
    story.append(header_bar)
    story.append(Spacer(1, 15))

    elements = parse_markdown(markdown_content)
    for elem in elements:
        if elem["type"] == "header":
            escaped = html.escape(elem["text"])
            if elem["level"] == 1:
                story.append(Paragraph(escaped, title_style))
            elif elem["level"] == 2:
                story.append(Paragraph(escaped, h2_style))
            else:
                story.append(Paragraph(escaped, h3_style))
        elif elem["type"] == "bullet":
            escaped = html.escape(elem["text"])
            b_style = get_bullet_style(elem["level"])
            story.append(Paragraph(f"• {escaped}", b_style))
        elif elem["type"] == "paragraph":
            escaped = html.escape(elem["text"])
            story.append(Paragraph(escaped, body_style))

    confidential_style = ParagraphStyle('CitiConf', parent=body_style, fontSize=8, textColor=colors.gray, alignment=1)
    story.append(Spacer(1, 25))
    story.append(Paragraph("CITI INTERNAL USE ONLY - STRICTLY CONFIDENTIAL", confidential_style))
    
    doc.build(story)

def export_report_to_pptx(markdown_content: str, output_path: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"PPTX FALLBACK - MOCK GENERATED\n\n{markdown_content}")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    citi_blue_rgb = RGBColor(0, 59, 112)
    citi_red_rgb = RGBColor(238, 49, 36)
    charcoal_rgb = RGBColor(34, 34, 34)

    # Parse markdown into elements
    elements = parse_markdown(markdown_content)

    # Parse elements into sections
    sections = []
    current_title = None
    current_content = []
    
    for elem in elements:
        if elem["type"] == "header":
            if current_title is not None or current_content:
                sections.append({
                    "title": current_title or "Report Details",
                    "content": current_content
                })
            current_title = elem["text"]
            current_content = []
        else:
            current_content.append(elem)
            
    if current_title is not None or current_content:
        sections.append({
            "title": current_title or "Report Details",
            "content": current_content
        })

    # Create Title Slide
    first_section_title = sections[0]["title"] if sections else "Citi Operational Risk Analytics Report"
    
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    
    shapes = slide.shapes
    header_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
    tf = header_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = first_section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = citi_blue_rgb

    # Accent Red line
    line_shape = shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(12.333), Inches(0.08))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = citi_red_rgb
    try:
        line_shape.line.color.rgb = citi_red_rgb
    except Exception:
        pass

    # Create Content Slides
    max_items_per_slide = 6
    for idx, sec in enumerate(sections):
        # Skip creating a content slide for the first section if it has no content
        if idx == 0 and not sec["content"]:
            continue
            
        title = sec["title"] or "Report Details"
        content_items = sec["content"]
        
        # Split content_items into chunks of max_items_per_slide
        chunks = [content_items[i:i + max_items_per_slide] for i in range(0, len(content_items), max_items_per_slide)]
        
        for chunk_idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[1]) 
            shapes = slide.shapes
            title_shape = shapes.title
            
            display_title = title
            if len(chunks) > 1:
                display_title += f" (Cont. {chunk_idx + 1})"
            title_shape.text = display_title
            title_shape.text_frame.paragraphs[0].font.color.rgb = citi_blue_rgb
            
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            first_para = True
            for elem in chunk:
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                
                p.text = elem["text"]
                p.font.color.rgb = charcoal_rgb
                
                if elem["type"] == "bullet":
                    p.level = elem["level"]
                else:
                    p.level = 0

    prs.save(output_path)
